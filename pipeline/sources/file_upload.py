"""Extract plain text from uploaded study files (.pdf, .pptx, .docx, .txt)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from docx import Document
import pdfplumber


def extract_txt_text(path: str) -> str:
    return Path(path).read_text(encoding="utf-8", errors="replace")


def extract_pdf_text(path: str) -> str:
    blocks: list[str] = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            if text.strip():
                blocks.append(text.strip())
    return "\n\n".join(blocks)


def extract_pptx_text(path: str) -> str:
    prs = Presentation(path)
    text_blocks: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = " ".join(run.text for run in para.runs)
                    if text.strip():
                        text_blocks.append(text.strip())
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                text_blocks.append(notes.strip())
    return "\n".join(text_blocks)


def extract_docx_text(path: str) -> str:
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_file_text(path: str) -> str:
    ext = Path(path).suffix.lower()
    if ext == ".txt":
        return extract_txt_text(path)
    if ext == ".pdf":
        return extract_pdf_text(path)
    if ext == ".pptx":
        return extract_pptx_text(path)
    if ext == ".docx":
        return extract_docx_text(path)
    raise ValueError(f"Unsupported file type: {ext}")
